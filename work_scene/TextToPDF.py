# -*- coding: utf-8 -*-
"""
@Name:        TextToPDF.py
@Describe:    XXX
@Time:        2025/01/20 15:16:01
@Author:      Ray
@Version:     1.0
"""
import csv
import sys
import json
from pptx.util import Inches, Pt
from pptx import Presentation
from weasyprint import HTML
from docx import Document
from openpyxl import Workbook


def txt_to_pdf(txt_file, pdf_file):
    # 读取并写入pdf文件
    with open(txt_file, 'r', encoding='utf-8') as f:
        content = f.read()

    html_content = f"<html><body><pre>{content}</pre></body></html>"
    HTML(string=html_content).write_pdf(pdf_file)


def txt_to_word(txt_file, word_file):
    # 读取并写入word文件
    document = Document()
    with open(txt_file, 'r', encoding='utf-8') as f:
        content = f.read()

    document.add_paragraph(content)
    document.save(word_file)


def txt_to_excel(txt_file, excel_file):
    wb = Workbook()
    ws = wb.active
    with open(txt_file, 'r', encoding='utf-8') as f:
        for line in f:
            ws.append([line.strip()])
    wb.save(excel_file)


def txt_to_html(txt_file, html_file):
    with open(txt_file, 'r', encoding='utf-8') as f:
        content = f.read()
    html_content = f"<html><body><pre>{content}</pre></body></html>"
    with open(html_file, 'w', encoding='utf-8') as f:
        f.write(html_content)


def txt_to_csv(txt_file, csv_file):
    with open(txt_file, 'r', encoding='utf-8') as f:
        lines = f.readlines()
    with open(csv_file, 'w', newline='', encoding='utf-8') as csvfile:
        writer = csv.writer(csvfile)
        for line in lines:
            writer.writerow([line.strip()])


def txt_to_json(txt_file, json_file):
    data = {}
    with open(txt_file, 'r', encoding='utf-8') as f:
        lines = f.readlines()
    for i in range(0, len(lines), 2):
        key = lines[i].strip()
        value = lines[i + 1].strip()
        data[key] = value
    with open(json_file, 'w', encoding='utf-8') as jf:
        json.dump(data, jf, ensure_ascii=False, indent=4)


def txt_to_ppt(txt_file, ppt_file):
    # 创建一个新的演示文稿
    prs = Presentation()

    # 选择一个不带标题和副标题占位符的布局
    blank_slide_layout = prs.slide_layouts[6]

    # 读取文本文件
    with open(txt_file, 'r', encoding='utf-8') as f:
        content = f.read()

    # 分割文本为合适的段落大小
    max_chars_per_slide = 500  # 可以根据需要调整每张幻灯片的最大字符数
    paragraphs = [content[i:i + max_chars_per_slide]
                  for i in range(0, len(content), max_chars_per_slide)]

    # 为每个段落创建一个新幻灯片
    for paragraph in paragraphs:
        slide = prs.slides.add_slide(blank_slide_layout)

        # 添加一个文本框
        left = top = Inches(1)
        width = Inches(8)
        height = Inches(5)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        # 设置文本框中的段落
        p = tf.add_paragraph()
        p.text = paragraph
        p.font.size = Pt(14)  # 设置字体大小

    # 保存演示文稿
    prs.save(ppt_file)


def main():
    txt_file = "work_scene/inputs/test.txt"
    dest_dir = "work_scene/outputs/"
    # 封装函数命令
    exchange_file = {
        'pdf': txt_to_pdf,
        'docx': txt_to_word,
        'xlsx': txt_to_excel,
        'html': txt_to_html,
        'csv': txt_to_csv,
        'json': txt_to_json,
        'pptx': txt_to_ppt,
    }
    extent = input("请输入要转化的文件后缀: ")
    if extent not in exchange_file.keys():
        print("请输入正确的扩展名")
        sys.exit()
    # 利用partition字符串截取，获取文件名
    name_file = txt_file.rpartition("/")[-1].partition(".")[0] + '.' + extent
    # 拼接文件路径
    dest_file = dest_dir + name_file
    exchange_file[extent](txt_file, dest_file)


if __name__ == '__main__':
    main()
